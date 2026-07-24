import sys
import os
import subprocess

# Self-healing dependency installer
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("[INFO] python-pptx not detected. Installing library...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

def create_deck():
    prs = Presentation()
    # Use widescreen layout (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Theme Colors matching At Your Service dashboard
    BG_COLOR = RGBColor(15, 23, 42)        # Slate 900 (#0f172a)
    TEXT_PRIMARY = RGBColor(248, 250, 252) # Slate 50 (#f8fafc)
    TEXT_MUTED = RGBColor(148, 163, 184)   # Slate 400 (#94a3b8)
    ACCENT_COLOR = RGBColor(99, 102, 241)  # Indigo (#6366f1)
    CYAN_COLOR = RGBColor(6, 182, 212)     # Cyan (#06b6d4)
    
    blank_layout = prs.slide_layouts[6] # Blank slide
    
    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Glassmorphic look)
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = BG_COLOR
    
    # Brand logo block
    logo_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf1 = logo_box.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    run_brand1 = p1.add_run()
    run_brand1.text = "AT YOUR "
    run_brand1.font.name = "Outfit"
    run_brand1.font.size = Pt(64)
    run_brand1.font.bold = True
    run_brand1.font.color.rgb = TEXT_PRIMARY
    
    run_brand2 = p1.add_run()
    run_brand2.text = "SERVICE"
    run_brand2.font.name = "Outfit"
    run_brand2.font.size = Pt(64)
    run_brand2.font.bold = True
    run_brand2.font.color.rgb = ACCENT_COLOR
    
    p2 = tf1.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(14)
    run_sub = p2.add_run()
    run_sub.text = "A Next-Generation Restaurant Analytics & Kitchen Operations Dashboard"
    run_sub.font.name = "Inter"
    run_sub.font.size = Pt(22)
    run_sub.font.color.rgb = TEXT_MUTED
    
    p3 = tf1.add_paragraph()
    p3.alignment = PP_ALIGN.LEFT
    p3.space_before = Pt(40)
    run_details = p3.add_run()
    run_details.text = "Project Presentation for Academic Review | Powered by Flask, MySQL & Vanilla JS"
    run_details.font.name = "Inter"
    run_details.font.size = Pt(14)
    run_details.font.italic = True
    run_details.font.color.rgb = CYAN_COLOR
    
    # ----------------------------------------------------
    # SLIDE 2: Problem Statement & Vision
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = BG_COLOR
    
    # Title
    t_box2 = slide2.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.0))
    tf2 = t_box2.text_frame
    p_t2 = tf2.paragraphs[0]
    r_t2 = p_t2.add_run()
    r_t2.text = "Problem Statement & Vision"
    r_t2.font.name = "Outfit"
    r_t2.font.size = Pt(36)
    r_t2.font.bold = True
    r_t2.font.color.rgb = ACCENT_COLOR
    
    # Content Columns (Left: Problem, Right: Vision)
    # Left Column (Problem)
    left_box = slide2.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    
    p_l1 = tf_l.paragraphs[0]
    r_l1 = p_l1.add_run()
    r_l1.text = "Current Challenges in Restaurant IT:"
    r_l1.font.name = "Outfit"
    r_l1.font.size = Pt(20)
    r_l1.font.bold = True
    r_l1.font.color.rgb = CYAN_COLOR
    
    bullets_l = [
        "Disconnected Systems: Front-of-House transaction systems (POS) rarely sync dynamically with Back-of-House kitchen lines.",
        "Resource Waste: Inventory management and staffing are reactive due to a lack of local daily forecasting systems.",
        "Rigid UI Design: Traditional POS terminal user interfaces are outdated, confusing, and lack responsive modern design patterns."
    ]
    for b in bullets_l:
        p_b = tf_l.add_paragraph()
        p_b.space_before = Pt(12)
        r_b = p_b.add_run()
        r_b.text = "•  " + b.split(":")[0] + ":"
        r_b.font.bold = True
        r_b.font.size = Pt(15)
        r_b.font.color.rgb = TEXT_PRIMARY
        r_b.font.name = "Inter"
        
        r_b_desc = p_b.add_run()
        r_b_desc.text = b.split(":")[1]
        r_b_desc.font.size = Pt(15)
        r_b_desc.font.color.rgb = TEXT_MUTED
        r_b_desc.font.name = "Inter"
        
    # Right Column (Vision)
    right_box = slide2.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    
    p_r1 = tf_r.paragraphs[0]
    r_r1 = p_r1.add_run()
    r_r1.text = "Our Unified Solution (At Your Service):"
    r_r1.font.name = "Outfit"
    r_r1.font.size = Pt(20)
    r_r1.font.bold = True
    r_r1.font.color.rgb = ACCENT_COLOR
    
    bullets_r = [
        "Direct Operations Bridge: Live sync between customer sales, an automated simulation feed, and a kitchen queue board.",
        "Embedded Math Intelligence: Python forecasting models that project future revenue, taking seasonality weights into account.",
        "Interactive Premium UX: Dark theme glassmorphism, responsive 3D card tilt gestures, and clean data visualizations."
    ]
    for b in bullets_r:
        p_b = tf_r.add_paragraph()
        p_b.space_before = Pt(12)
        r_b = p_b.add_run()
        r_b.text = "•  " + b.split(":")[0] + ":"
        r_b.font.bold = True
        r_b.font.size = Pt(15)
        r_b.font.color.rgb = TEXT_PRIMARY
        r_b.font.name = "Inter"
        
        r_b_desc = p_b.add_run()
        r_b_desc.text = b.split(":")[1]
        r_b_desc.font.size = Pt(15)
        r_b_desc.font.color.rgb = TEXT_MUTED
        r_b_desc.font.name = "Inter"
        
    # ----------------------------------------------------
    # SLIDE 3: System Architecture
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    slide3.background.fill.solid()
    slide3.background.fill.fore_color.rgb = BG_COLOR
    
    # Title
    t_box3 = slide3.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.0))
    tf3 = t_box3.text_frame
    p_t3 = tf3.paragraphs[0]
    r_t3 = p_t3.add_run()
    r_t3.text = "System Architecture & Data Flow"
    r_t3.font.name = "Outfit"
    r_t3.font.size = Pt(36)
    r_t3.font.bold = True
    r_t3.font.color.rgb = ACCENT_COLOR
    
    # Details Box (Left: Architecture Stack, Right: Pipeline Sequence)
    arch_left = slide3.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_al = arch_left.text_frame
    tf_al.word_wrap = True
    
    p_al = tf_al.paragraphs[0]
    r_al = p_al.add_run()
    r_al.text = "The Three-Tier Architecture:"
    r_al.font.name = "Outfit"
    r_al.font.size = Pt(20)
    r_al.font.bold = True
    r_al.font.color.rgb = CYAN_COLOR
    
    tiers = [
        "Database Layer: MySQL Relational Server. Holds the sales table with structural indexes configured for rapid datetime filtering and status sorting.",
        "Application Layer: Python Flask Server. Exposes CRUD APIs, maps dates, serves static assets, runs linear forecasting, and handles background simulator threads.",
        "Presentation Layer: SPA Frontend. Styled with HSL CSS variables, utilizing Chart.js for visualization, html2pdf.js for receipts, and custom JS event loops."
    ]
    for t in tiers:
        p_t = tf_al.add_paragraph()
        p_t.space_before = Pt(12)
        r_t = p_t.add_run()
        r_t.text = "•  " + t.split(":")[0] + ":"
        r_t.font.bold = True
        r_t.font.size = Pt(15)
        r_t.font.color.rgb = TEXT_PRIMARY
        
        r_t_desc = p_t.add_run()
        r_t_desc.text = t.split(":")[1]
        r_t_desc.font.size = Pt(14)
        r_t_desc.font.color.rgb = TEXT_MUTED
        
    # Process Flow (Right Column)
    arch_right = slide3.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_ar = arch_right.text_frame
    tf_ar.word_wrap = True
    
    p_ar = tf_ar.paragraphs[0]
    r_ar = p_ar.add_run()
    r_ar.text = "Lifecycle of a Transaction Order:"
    r_ar.font.name = "Outfit"
    r_ar.font.size = Pt(20)
    r_ar.font.bold = True
    r_ar.font.color.rgb = ACCENT_COLOR
    
    steps = [
        "1. Ingestion: Simulator daemon (or manual order) makes a POST request to /api/sales, writing the transaction directly to MySQL with status 'PENDING'.",
        "2. KDS Broadcast: Front-of-house KDS tab polls /api/kds, captures the pending order, groups items, and binds it to a styled card with active timers.",
        "3. Preparation Update: Kitchen staff click preparation buttons to progress tickets (Pending ➔ Preparing ➔ Ready ➔ Completed) via POST updates.",
        "4. PDF Document Compile: Client retrieves item data on-click, styles an HTML document in memory, applies 5% GST, and exports an A4 PDF receipt."
    ]
    for s in steps:
        p_s = tf_ar.add_paragraph()
        p_s.space_before = Pt(10)
        r_s = p_s.add_run()
        r_s.text = s
        r_s.font.size = Pt(14)
        r_s.font.color.rgb = TEXT_MUTED
        r_s.font.name = "Inter"
        
    # ----------------------------------------------------
    # SLIDE 4: Front of House Features
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    slide4.background.fill.solid()
    slide4.background.fill.fore_color.rgb = BG_COLOR
    
    t_box4 = slide4.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.0))
    tf4 = t_box4.text_frame
    p_t4 = tf4.paragraphs[0]
    r_t4 = p_t4.add_run()
    r_t4.text = "Key Features: Front of House & Ledger"
    r_t4.font.name = "Outfit"
    r_t4.font.size = Pt(36)
    r_t4.font.bold = True
    r_t4.font.color.rgb = ACCENT_COLOR
    
    content4 = slide4.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(4.8))
    tf_c4 = content4.text_frame
    tf_c4.word_wrap = True
    
    feats_foh = [
        "Interactive KPI Summary: Tracks Total Revenue, total orders placed, and Average Order Value (AOV) dynamically in Indian Rupees (₹).",
        "Transaction Audit Ledger: A fully paginated table showing order dates, unit items, quantity multipliers, payment modes, and order types.",
        "Client-Side PDF Receipt Compiler: Generates fully-styled, professional A4 PDF invoices directly in the browser using html2pdf.js, calculating 5% GST and generating automatic downloads, saving server processing power.",
        "Real-Time Sales Simulator: Multi-threaded background process in Python writing mock orders every 2-5 seconds, creating live database changes.",
        "Dynamic 3D Hover Tilt Gestures: Applies CSS perspective transforms and translateZ translations on KPI cards so they tilt dynamically relative to mouse movement, creating a spectacular parallax holographic effect."
    ]
    for f in feats_foh:
        p_f = tf_c4.add_paragraph()
        if tf_c4.paragraphs[0] == p_f:
            p_f.space_before = Pt(0)
        else:
            p_f.space_before = Pt(14)
            
        r_f = p_f.add_run()
        r_f.text = "•  " + f.split(":")[0] + ":"
        r_f.font.bold = True
        r_f.font.size = Pt(16)
        r_f.font.color.rgb = TEXT_PRIMARY
        
        r_f_desc = p_f.add_run()
        r_f_desc.text = f.split(":")[1]
        r_f_desc.font.size = Pt(16)
        r_f_desc.font.color.rgb = TEXT_MUTED
        
    # ----------------------------------------------------
    # SLIDE 5: Back of House Features
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    slide5.background.fill.solid()
    slide5.background.fill.fore_color.rgb = BG_COLOR
    
    t_box5 = slide5.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.0))
    tf5 = t_box5.text_frame
    p_t5 = tf5.paragraphs[0]
    r_t5 = p_t5.add_run()
    r_t5.text = "Key Features: Kitchen Display System (KDS)"
    r_t5.font.name = "Outfit"
    r_t5.font.size = Pt(36)
    r_t5.font.bold = True
    r_t5.font.color.rgb = ACCENT_COLOR
    
    content5 = slide5.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(4.8))
    tf_c5 = content5.text_frame
    tf_c5.word_wrap = True
    
    feats_boh = [
        "Live Kitchen Queue: Visualizes active orders on a card grid, automatically hiding completed records once served.",
        "Visual 3D Category Headers: Maps category icons and 3D mockups at the top of cards (Pizza for Mains, Burger for Appetizers, Drinks for Beverages) so chefs can parse order types at a single glance.",
        "Interactive Pipeline Controls: Statuses transition seamlessly (Pending ➔ Preparing ➔ Ready ➔ Completed) directly from the card footer.",
        "Elapsed Ticket Timers: Displays preparation timers ticking up in real-time (ticking every 1 second) showing precisely how long an order has been active.",
        "Critical Warnings: Visual alert system changes timer text to glowing red and flashes when tickets remain un-served for more than 10 minutes."
    ]
    for f in feats_boh:
        p_f = tf_c5.add_paragraph()
        if tf_c5.paragraphs[0] == p_f:
            p_f.space_before = Pt(0)
        else:
            p_f.space_before = Pt(14)
            
        r_f = p_f.add_run()
        r_f.text = "•  " + f.split(":")[0] + ":"
        r_f.font.bold = True
        r_f.font.size = Pt(16)
        r_f.font.color.rgb = TEXT_PRIMARY
        
        r_f_desc = p_f.add_run()
        r_f_desc.text = f.split(":")[1]
        r_f_desc.font.size = Pt(16)
        r_f_desc.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 6: Predictive Forecasting & Math Model
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    slide6.background.fill.solid()
    slide6.background.fill.fore_color.rgb = BG_COLOR
    
    t_box6 = slide6.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.0))
    tf6 = t_box6.text_frame
    p_t6 = tf6.paragraphs[0]
    r_t6 = p_t6.add_run()
    r_t6.text = "Predictive Forecasting & Mathematical Model"
    r_t6.font.name = "Outfit"
    r_t6.font.size = Pt(36)
    r_t6.font.bold = True
    r_t6.font.color.rgb = ACCENT_COLOR
    
    content6 = slide6.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(4.8))
    tf_c6 = content6.text_frame
    tf_c6.word_wrap = True
    
    model_points = [
        "Additive Trend + Day-of-Week Seasonality Model: Pure mathematical algorithm written in standard Python, bypassing heavy ML library overheads.",
        "Linear Regression (y = mx + b): Calculated dynamically over historical sales to isolate general revenue velocity and baseline direction.",
        "Day-of-Week Seasonality Weighting: Computes normal weekday factor averages relative to overall averages (e.g. Sunday weights peak at 1.4x baseline, while Tuesdays sit at 0.8x).",
        "Forecasting Equations:\n  - Trend Projection: T(t) = m * t + b\n  - Seasonal Value: F(t) = T(t) * Seasonal_Factor(Weekday(t))\n  - Shaded Confidence Band (90%): Lower Bound = F(t) * 0.90 | Upper Bound = F(t) * 1.10",
        "Confidence Shading: Rendered dynamically in Chart.js with transparent line fills mapping future projection limits."
    ]
    for mp in model_points:
        p_mp = tf_c6.add_paragraph()
        if tf_c6.paragraphs[0] == p_mp:
            p_mp.space_before = Pt(0)
        else:
            p_mp.space_before = Pt(12)
            
        r_mp = p_mp.add_run()
        r_mp.text = "•  " + mp.split(":")[0] + ":"
        r_mp.font.bold = True
        r_mp.font.size = Pt(15)
        r_mp.font.color.rgb = TEXT_PRIMARY
        
        r_mp_desc = p_mp.add_run()
        r_mp_desc.text = mp.split(":")[1]
        r_mp_desc.font.size = Pt(14)
        r_mp_desc.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 7: Rebranding & Aesthetics
    # ----------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    slide7.background.fill.solid()
    slide7.background.fill.fore_color.rgb = BG_COLOR
    
    t_box7 = slide7.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.0))
    tf7 = t_box7.text_frame
    p_t7 = tf7.paragraphs[0]
    r_t7 = p_t7.add_run()
    r_t7.text = "Rebranding, Design & Security Gate"
    r_t7.font.name = "Outfit"
    r_t7.font.size = Pt(36)
    r_t7.font.bold = True
    r_t7.font.color.rgb = ACCENT_COLOR
    
    content7 = slide7.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(4.8))
    tf_c7 = content7.text_frame
    tf_c7.word_wrap = True
    
    design_points = [
        "Login Security Portal: Full-screen entry gate restricting app container visibility. Employs localStorage session tracking.",
        "Visual Dish Collage Background: Sets an octane food collage (pizza, burgers, desserts, drinks) beneath a dimming glassmorphic overlay for a highly premium aesthetic.",
        "At Your Service Rebranding: Full user-facing rename (sidebar headers, title tags, PDF templates) while retaining backward database compatibility.",
        "Custom 3D Graphics: Leverages 3D isometric octane renders for KDS covers and New Order modal category previews, updating automatically as users switch categories."
    ]
    for dp in design_points:
        p_dp = tf_c7.add_paragraph()
        if tf_c7.paragraphs[0] == p_dp:
            p_dp.space_before = Pt(0)
        else:
            p_dp.space_before = Pt(14)
            
        r_dp = p_dp.add_run()
        r_dp.text = "•  " + dp.split(":")[0] + ":"
        r_dp.font.bold = True
        r_dp.font.size = Pt(16)
        r_dp.font.color.rgb = TEXT_PRIMARY
        
        r_dp_desc = p_dp.add_run()
        r_dp_desc.text = dp.split(":")[1]
        r_dp_desc.font.size = Pt(16)
        r_dp_desc.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 8: Tech Stack summary
    # ----------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    slide8.background.fill.solid()
    slide8.background.fill.fore_color.rgb = BG_COLOR
    
    t_box8 = slide8.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.0))
    tf8 = t_box8.text_frame
    p_t8 = tf8.paragraphs[0]
    r_t8 = p_t8.add_run()
    r_t8.text = "Technology Stack & Core Components"
    r_t8.font.name = "Outfit"
    r_t8.font.size = Pt(36)
    r_t8.font.bold = True
    r_t8.font.color.rgb = ACCENT_COLOR
    
    content8 = slide8.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(4.8))
    tf_c8 = content8.text_frame
    tf_c8.word_wrap = True
    
    techs = [
        "Backend Framework: Python 3.14 (with Flask micro-routing for API endpoints).",
        "Database Engine: MySQL Server (local instance, port 3306, user root) managed in MySQL Workbench.",
        "Interactive Charts: Chart.js (drawing linear trends, category doughnuts, hourly peak curves, and forecast confidence lines).",
        "PDF Invoices: html2pdf.js (CDN bundle, compiling and rendering HTML receipt blocks to PDF on the client side).",
        "Aesthetic Styling: CSS Custom HSL variables, flexbox grid, and dark/light glassmorphism layouts."
    ]
    for t in techs:
        p_t = tf_c8.add_paragraph()
        if tf_c8.paragraphs[0] == p_t:
            p_t.space_before = Pt(0)
        else:
            p_t.space_before = Pt(14)
            
        r_t = p_t.add_run()
        r_t.text = "•  " + t.split(":")[0] + ":"
        r_t.font.bold = True
        r_t.font.size = Pt(16)
        r_t.font.color.rgb = TEXT_PRIMARY
        
        r_t_desc = p_t.add_run()
        r_t_desc.text = t.split(":")[1]
        r_t_desc.font.size = Pt(16)
        r_t_desc.font.color.rgb = TEXT_MUTED

    # Save presentation
    filename = "At_Your_Service_Presentation.pptx"
    prs.save(filename)
    print(f"Presentation saved successfully as '{filename}'.")

if __name__ == "__main__":
    create_deck()
